import csv
import json

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.rag.loaders import load_document


def test_common_document_loaders(tmp_path):
    (tmp_path / "sample.txt").write_text("plain text", encoding="utf-8")
    (tmp_path / "sample.md").write_text("# markdown", encoding="utf-8")
    (tmp_path / "sample.html").write_text("<h1>HTML</h1>", encoding="utf-8")
    (tmp_path / "sample.json").write_text(
        json.dumps({"type": "json"}), encoding="utf-8"
    )

    with (tmp_path / "sample.csv").open("w", newline="") as file:
        csv.writer(file).writerows([["name", "value"], ["HY-chat", "RAG"]])

    document = Document()
    document.add_paragraph("Word document")
    document.save(tmp_path / "sample.docx")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "PowerPoint"
    presentation.save(tmp_path / "sample.pptx")

    workbook = Workbook()
    workbook.active["A1"] = "Excel"
    workbook.save(tmp_path / "sample.xlsx")

    for path in tmp_path.iterdir():
        sections = load_document(path)
        assert sections
        assert any(section.text.strip() for section in sections)
