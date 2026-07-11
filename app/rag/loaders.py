from __future__ import annotations

import csv
import json
from dataclasses import dataclass, field
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from app.core.types import JsonObject

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".txt",
    ".md",
    ".html",
    ".htm",
    ".csv",
    ".json",
}


@dataclass
class LoadedSection:
    text: str
    metadata: JsonObject = field(default_factory=dict)


def _load_pdf(path: Path) -> list[LoadedSection]:
    reader = PdfReader(path)
    sections: list[LoadedSection] = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            sections.append(LoadedSection(text=text, metadata={"page": index}))
    return sections


def _load_docx(path: Path) -> list[LoadedSection]:
    document = DocxDocument(path)
    text = "\n".join(
        paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()
    )
    return [LoadedSection(text=text)] if text else []


def _load_pptx(path: Path) -> list[LoadedSection]:
    presentation = Presentation(path)
    sections: list[LoadedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            sections.append(
                LoadedSection(text="\n".join(texts), metadata={"slide": index})
            )
    return sections


def _load_xlsx(path: Path) -> list[LoadedSection]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections: list[LoadedSection] = []
    for sheet in workbook.worksheets:
        rows = [
            "\t".join("" if value is None else str(value) for value in row)
            for row in sheet.iter_rows(values_only=True)
        ]
        text = "\n".join(row for row in rows if row.strip())
        if text:
            sections.append(LoadedSection(text=text, metadata={"sheet": sheet.title}))
    workbook.close()
    return sections


def _load_csv(path: Path) -> list[LoadedSection]:
    with path.open("r", encoding="utf-8-sig", errors="ignore", newline="") as file:
        rows = ["\t".join(row) for row in csv.reader(file)]
    text = "\n".join(row for row in rows if row.strip())
    return [LoadedSection(text=text)] if text else []


def _load_json(path: Path) -> list[LoadedSection]:
    data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    return [LoadedSection(text=json.dumps(data, ensure_ascii=False, indent=2))]


def _load_html(path: Path) -> list[LoadedSection]:
    html = path.read_text(encoding="utf-8", errors="ignore")
    text = BeautifulSoup(html, "html.parser").get_text("\n", strip=True)
    return [LoadedSection(text=text)] if text else []


def load_document(path: str | Path) -> list[LoadedSection]:
    file_path = Path(path)
    extension = file_path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(
            f"Unsupported document type '{extension}'. Supported: {supported}"
        )

    if extension == ".pdf":
        sections = _load_pdf(file_path)
    elif extension == ".docx":
        sections = _load_docx(file_path)
    elif extension == ".pptx":
        sections = _load_pptx(file_path)
    elif extension == ".xlsx":
        sections = _load_xlsx(file_path)
    elif extension == ".csv":
        sections = _load_csv(file_path)
    elif extension == ".json":
        sections = _load_json(file_path)
    elif extension in {".html", ".htm"}:
        sections = _load_html(file_path)
    else:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        sections = [LoadedSection(text=text)] if text.strip() else []

    if not sections:
        raise ValueError("No extractable text was found in the document")
    return sections
